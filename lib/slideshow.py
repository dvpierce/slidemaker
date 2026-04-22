import os
import gc
import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from threading import Lock

from lib.helpers import dedup, imghandler, my_logger
import yaml
import datetime

class slideshow:
    no_transition = """<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
    <mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" Requires="p14">
      <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" p14:dur="0" advClick="0" advTm="{timedelay}"/>
    </mc:Choice>
    <mc:Fallback>
      <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" advClick="0" advTm="{timedelay}"/>
    </mc:Fallback>
  </mc:AlternateContent>"""

    fade_transition = """<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
    <mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" Requires="p14">
      <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" p14:dur="700" advClick="0" advTm="{timedelay}">
        <p:fade/>
      </p:transition>
    </mc:Choice>
    <mc:Fallback>
      <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advClick="0" advTm="{timedelay}">
        <p:fade/>
      </p:transition>
    </mc:Fallback>
  </mc:AlternateContent>"""

    wipe_transition = """<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow" advClick="0" advTm="{timedelay}">
        <p:wipe/>
    </p:transition>"""

    push_transition = """<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow" advClick="0" advTm="{timedelay}">
        <p:push dir="u"/>
    </p:transition>"""

    ripple_transition = """<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
    <mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" Requires="p14">
        <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow" p14:dur="1400" advClick="0" advTm="{timedelay}">
            <p14:ripple/>
        </p:transition>
    </mc:Choice>
    <mc:Fallback>
        <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow" advClick="0" advTm="{timedelay}">
            <p:fade/>
        </p:transition>
    </mc:Fallback>
</mc:AlternateContent>"""

    transitions = [
        {"name": "fade", "duration": 700},
        {"name": "ripple", "duration": 2000},
    ]

    def __init__(self,
                 input_dir=None,
                 output_file=None,
                 overwrite=False,
                 slide_w=13.333,
                 slide_h=7.5,
                 transition="none",
                 auto_adjust=False,
                 bgcolor="ffffff",
                 slide_duration_sec=5,
                 blurbg=None,
                 deduplicate=False,
                 duplicate_threshold=12,
                 resample=0,
                 image_quality=75,
                 logfile=None,
                 subfolders=False,
                 titles=False,
                 captions_file=None):
        self.input_dir = input_dir
        self.output_file = output_file
        self.overwrite = overwrite
        self.slide_w = slide_w
        self.slide_h = slide_h
        self.transition = transition
        self.auto_adjust = auto_adjust
        self.bgcolor = bgcolor
        self.slide_duration_sec = slide_duration_sec
        self.blurbg = blurbg
        self.deduplicate = deduplicate
        self.duplicate_threshold = duplicate_threshold
        self.resample=resample
        self.subfolders = subfolders
        self.titles = titles
        self.captions = captions_file
        self.image_quality = image_quality
        if self.captions:
            self.enable_captions = True
            self.captiondata = yaml.safe_load(open(self.captions, "r").read())
        self.showstopper = Lock()
        self.out_log = my_logger(file=logfile)

        if not sys.maxsize > 2**32:
            print("You are using a 32-bit version of Python: you may encounter issues with photo sets exceeding 4GB.")
        return

    def output(self, *args, **kwargs):
        self.out_log.logit(*args, **kwargs)

    def _insert_caption(self):
        if self.enable_captions:
            if self.current_imagehandler.filename not in self.captiondata.keys():
                return
            else:
                left = Inches(0)
                top = Inches(self.slide_h - 0.5)
                width = Inches(self.slide_w)
                height = Inches(0.5)

                txBox = self.current_slide.shapes.add_textbox(left, top, width, height)
                txBox.fill.solid()

                r, g, b = imghandler.convert_RGB(self.bgcolor)
                txBox.fill.fore_color.rgb = RGBColor(r, g, b)

                txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = txBox.text_frame

                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = self.captiondata[self.current_imagehandler.filename]

                txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                txBox.text_frame.fit_text(font_family='Calibri', max_size=24)

                # After a fit_text operation, reset the location/dimensions of the text box so it still covers the bottom half inch of the slide.
                txBox.width = width
                txBox.height = height
                txBox.left = left
                txBox.top = top

                tf.word_wrap = False
                font = run.font
                font.name = "Garamond"
                font.bold = False
                font.italic = True
                r, g, b = imghandler.get_appropriate_text_color(self.bgcolor)
                font.color.rgb = RGBColor(r, g, b)

        return

    def _insert_title_slide(self, title_text):
        self.current_slide_layout = self.current_presentation.slide_layouts[6]
        self.current_slide = self.current_presentation.slides.add_slide(self.current_slide_layout)
        self.current_slide.background.fill.solid()
        r, g, b = imghandler.convert_RGB(self.bgcolor)
        self.current_slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

        left = Inches(1)
        top = Inches(self.slide_h/2 - 1)
        width = Inches(self.slide_w - 2)
        height = Inches(1.5)

        txBox = self.current_slide.shapes.add_textbox(left, top, width, height)
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = RGBColor(r, g, b)

        txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = txBox.text_frame
        tf.word_wrap = False

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        font = run.font
        font.name = 'Garamond'
        font.bold = True
        r, g, b = imghandler.get_appropriate_text_color(self.bgcolor)
        font.color.rgb = RGBColor(r, g, b)
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        txBox.text_frame.fit_text(font_family='Calibri', max_size=48)

    def _insert_slides(self, image_files):
        count = 0
        for img_path in image_files:
            count += 1
            gc.collect()
            self.output(f"Creating slide {count} of {len(image_files)}...", end="\r", flush=True)

            # Use a blank slide layout (index 6 is usually blank)
            self.current_slide_layout = self.current_presentation.slide_layouts[6]
            self.current_slide = self.current_presentation.slides.add_slide(self.current_slide_layout)
            self.current_slide.background.fill.solid()

            # Determine desired size and position of image.
            try:
                self.current_imagehandler = imghandler(img_path, slide_w=self.slide_w, slide_h=self.slide_h, image_quality=self.image_quality, maxres=self.resample)
            except:
                # If there's an error creating the image handler, skip this image. It's probably just a bad file.
                continue
            width, height = self.current_imagehandler.get_limits()
            xoffset, yoffset = self.current_imagehandler.get_offsets()

            # Insert background on slide.
            if self.blurbg:
                self.current_slide.shapes.add_picture(i.blur_stretch(), 0, 0, width=Inches(self.slide_w), height=Inches(self.slide_h))
            else:
                r, g, b = imghandler.convert_RGB(self.bgcolor)
                self.current_slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

            if self.resample != 0:
                self.current_imagehandler.resample()

            if self.auto_adjust:
                self.current_slide.shapes.add_picture(self.current_imagehandler.get_autoadjusted_image(), Inches(xoffset), Inches(yoffset), width=Inches(width), height=Inches(height))
            else:
                self.current_slide.shapes.add_picture(self.current_imagehandler.get_image(), Inches(xoffset), Inches(yoffset), width=Inches(width), height=Inches(height))

            self._insert_caption()

            # Add slide transition (or no transition.)
            if self.transition == "fade":
                transition = self.fade_transition
            elif self.transition == "wipe":
                transition = self.wipe_transition
            elif self.transition == "push":
                transition = self.push_transition
            elif self.transition == "ripple":
                transition = self.ripple_transition
            else:
                transition = self.no_transition

            xml = transition.format(timedelay = str(self.slide_duration_sec * 1000))
            fragment = parse_xml(xml)
            self.current_slide.element.insert(-1, fragment)

    def _find_subdirs(self):
        abs_root = os.path.abspath(self.input_dir)
        dir_paths = []

        for root, dirs, files in os.walk(abs_root):
            dir_paths.append(root)
        return dir_paths

    def create_image_slideshow(self):
        with self.showstopper:
            if os.path.exists(self.output_file) and not self.overwrite:
                raise FileExistsError(f"{output_file} already exists. Not set to overwrite.")

            self.current_presentation = Presentation()
            self.current_presentation.slide_width = Inches(self.slide_w)
            self.current_presentation.slide_height = Inches(self.slide_h)

            if self.subfolders:
                directories = self._find_subdirs()
            else:
                directories = [self.input_dir]

            dircount = 0
            for directory in sorted(directories):
                dircount += 1
                self.output(f"Processing folder {dircount} of {len(directories)}...")
                if self.titles:
                    # Insert title slide (using name of directory as the title text)
                    title_text = os.path.basename(directory)
                    self._insert_title_slide(title_text)
                if self.deduplicate:
                    # Deduplicate the contents of the folder
                    deduper = dedup(directory=directory, hash_size=8, hamming_diff=self.duplicate_threshold, out_log=self.out_log)
                    image_files = deduper.get_deduplicated_file_list()
                else:
                    # Or not
                    image_files = [f"{directory}{os.sep}{file}" for file in os.listdir(directory) if imghandler.is_supported(file)]

                # Insert slides - create one slide per photo.
                self._insert_slides(image_files)

            self.output("")
            self.output(f"Saving {self.output_file}...")
            if os.path.exists(self.output_file):
                self.output(f"Warning: {self.output_file} exists and will be overwritten. (--overwrite enabled.)")
            self.current_presentation.save(self.output_file)
            self.output(f"Successfully created: {self.output_file}")
        return True
