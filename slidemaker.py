import os
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image
import numpy as np
import argparse
import re

parser = argparse.ArgumentParser()
parser.add_argument("--width", help="Slide width, in inches", type=float, required=False, default=13.333)
parser.add_argument("--height", help="Slide height, in inches", type=float, required=False, default=7.5)
parser.add_argument("--input_dir", help="Path to files", type=str, required=False, default="./")
parser.add_argument("--output_file", help="Name of output file (pptx)", type=str, required=False, default="Presentation.pptx")
parser.add_argument("--duration", help="Duration for each slide.", type=float, required=False, default=5)
parser.add_argument("--overwrite", help="Overwrite existing files?", action="store_true")
parser.add_argument("--bgcolor", help="RGB code for background color. Default is black.", required=False, default="000000")
parser.add_argument("--threshold", help="Hamming distance for duplicate detection. Default is 12. Smaller numbers are less likely to detect duplicates, larger numbers are more likely to get false positives.", required=False, type=int, default=12)
args = parser.parse_args()

slide_w = args.width
slide_h = args.height


def get_image_fingerprint(image_path, hash_size=8):
    # 1. Open image and convert to grayscale ('L' mode)
    with Image.open(image_path) as img:
        img = img.convert('L')

        # 2. Resize to (hash_size + 1, hash_size)
        # We need an extra column to compare left vs. right pixels
        img = img.resize((hash_size + 1, hash_size), Image.Resampling.LANCZOS)

        # 3. Convert to a numpy array for easy comparison
        pixels = np.array(img)

        # 4. Compute differences between adjacent horizontal pixels
        # This creates a boolean array where True = left is brighter than right
        diff = pixels[:, 1:] > pixels[:, :-1]

        # 5. Convert the boolean array into a hexadecimal fingerprint
        decimal_value = 0
        hex_string = []
        for index, value in enumerate(diff.flatten()):
            if value:
                decimal_value += 2**(index % 8)
            if (index % 8) == 7:
                hex_string.append(hex(decimal_value)[2:].rjust(2, '0'))
                decimal_value = 0

        return "".join(hex_string)


def hamming_distance(hex_hash1, hex_hash2):
    # 1. Convert hex strings to integers
    val1 = int(hex_hash1, 16)
    val2 = int(hex_hash2, 16)

    # 2. XOR the two values
    # Bits that are different will result in a 1
    xor_result = val1 ^ val2

    # 3. Count the number of set bits (1s) in the binary result
    return bin(xor_result).count('1')


def getSize(image_file):
    w, h = Image.open(image_file).size
    return w, h


def getlimits(image_file):
    w, h = getSize(image_file)
    scalefactor = max([(w/slide_w), (h/slide_h)])
    return (w/scalefactor), (h/scalefactor)


def get_offsets(width, height):
    xoffset = (slide_w - width)/2
    yoffset = (slide_h - height)/2
    return xoffset, yoffset


def getRGB(hexstring):
    pattern = r'^#?([A-Fa-f0-9]{6}|[A-Fa-f0-9]{3})$'
    if not bool(re.match(pattern, hexstring)):
        raise Exception(f"{hexstring} is not a valid RGB color code.")
    hexstring = hexstring.replace("#", "").lower()
    if len(hexstring) == 3:
        r = int(hexstring[0]+hexstring[0], 16)
        g = int(hexstring[1]+hexstring[1], 16)
        b = int(hexstring[2]+hexstring[2], 16)
    elif len(hexstring) == 6:
        r = int(hexstring[0:2], 16)
        g = int(hexstring[2:4], 16)
        b = int(hexstring[4:6], 16)
    else:
        raise Exception(f"{hexstring} is not a valid RGB color code.")
    return RGBColor(r, g, b)


def find_dupes(image_files, duplicate_threshold=12):
    print("Fingerprinting images...")
    image_fingerprints = dict()
    count = 0
    for img_name in image_files:
        count += 1
        image_fingerprints[img_name] = get_image_fingerprint(img_name, hash_size=8)
        print(f"Processed {count} of {len(image_files)} images...", end="\r", flush=True)
    print("")

    print("Scanning for duplicates...")
    potential_duplicates = list()
    count = 0
    for img_name in image_files:
        count += 1
        this_image_hash = image_fingerprints[img_name]
        print(f"Checking image {count} of {len(image_files)} against {len(image_fingerprints.keys())-1} fingerprints", end="\r", flush=True)
        for other_file in image_fingerprints.keys():
            if other_file == img_name:
                continue
            if hamming_distance(this_image_hash, image_fingerprints[other_file]) <= duplicate_threshold:
                potential_duplicates.append(tuple(sorted((other_file, img_name))))
    print("")
    potential_duplicates = list(set(potential_duplicates))
    maxlength = max([ max([ len(x[0]), len(x[1]) ]) for x in potential_duplicates ])
    for x in potential_duplicates:
        print("These images are suspiciously similar:", x[0].ljust(maxlength), x[1].ljust(maxlength))
    return


def create_image_slideshow(input_dir=None,
                           output_file=None,
                           slide_duration_sec=5,
                           overwrite=False,
                           bgcolor="000000",
                           duplicate_threshold=12):
    # Initialize presentation
    prs = Presentation()

    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)

    valid_extensions = ('.avif', '.bmp', '.gif', '.j2k', '.jp2', '.jpx', '.pcx', '.tiff', '.tif', '.jpg', '.jpeg', '.png', '.webp')
    image_files = [f for f in os.listdir(input_dir) if f.lower().endswith(valid_extensions)]
    image_files = [ os.path.join(input_dir, img_name) for img_name in image_files ]
    image_files.sort()

    if not image_files:
        print("No images found in the specified folder.")
        return

    find_dupes(image_files, duplicate_threshold=duplicate_threshold)

    proceed = input("Would you like to proceed with creating the slideshow? [y/n]: ")
    if proceed.lower() not in ("y", "yes", "ok", "yup", "sure"):
        print("User cancelled.")
        return

    count = 0
    for img_path in image_files:
        count += 1
        print(f"Creating slide {count} of {len(image_files)}...", end="\r", flush=True)

        # Use a blank slide layout (index 6 is usually blank)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = getRGB(bgcolor)

        width, height = getlimits(img_path)
        xoffset, yoffset = get_offsets(width, height)

        # Add and resize image to fit the 10x7.5 slide exactly
        slide.shapes.add_picture(img_path, Inches(xoffset), Inches(yoffset), width=Inches(width), height=Inches(height))

        # Set timing: Access underlying XML to set 'Advance After' time
        # 'advTm' is in milliseconds
        slide_element = slide.element
        slide_element.set('advClick', '0')  # Don't wait for click
        slide_element.set('advTm', str(slide_duration_sec * 1000))

    print("")
    print(f"Saving {output_file}...")
    try:
        if os.path.exists(output_file) and (not overwrite):
            exit(f"{output_file} exists. Will not overwrite. Exiting.")
        else:
            if os.path.exists(output_file):
                print(f"Warning: {output_file} exists. (--overwrite selected.)")
            prs.save(output_file)
            print(f"Successfully created: {output_file}")
    except Exception as e:
        print("Error writing output file.")
        raise e


if __name__ == "__main__":
    create_image_slideshow(input_dir=args.input_dir,
                           output_file=args.output_file,
                           slide_duration_sec=args.duration,
                           overwrite=args.overwrite,
                           bgcolor=args.bgcolor,
                           duplicate_threshold=args.threshold)
